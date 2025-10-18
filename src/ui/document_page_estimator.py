"""
Document Page Count Estimator - Pre-extraction page/size estimation for cost calculations

Provides lightweight page count estimates WITHOUT calling paid document extractors.
Uses format-aware metadata reads where possible, with safe fallbacks.

Usage:
    from src.ui.document_page_estimator import estimate_document_pages

    page_count, confidence = estimate_document_pages(uploaded_file)
    # Returns: (15, 'high') for PDFs with readable metadata
    # Returns: (10, 'medium') for file-size based estimates
"""

import os
import logging
from pathlib import Path
from typing import Tuple, Union, BinaryIO

logger = logging.getLogger(__name__)

# Default page count when metadata unavailable (configurable via env)
DEFAULT_PAGE_ESTIMATE = int(os.getenv('DEFAULT_PAGE_ESTIMATE', '15'))

# File size heuristics (bytes per page estimates)
SIZE_PER_PAGE_TEXT_PDF = 30_000      # ~30KB per page for text-heavy PDFs
SIZE_PER_PAGE_SCANNED_PDF = 300_000  # ~300KB per page for scanned PDFs
SIZE_PER_PAGE_DOCX = 20_000          # ~20KB per page for DOCX
SIZE_PER_PAGE_GENERIC = 50_000       # ~50KB per page generic fallback


def estimate_document_pages(
    file_path: Union[Path, str, BinaryIO],
    filename: str = None
) -> Tuple[int, str]:
    """
    Estimate document page count without extraction.

    Tries format-specific metadata reads first, falls back to file size heuristics.

    Args:
        file_path: Path to document, file path string, or BytesIO file-like object
        filename: Original filename (required for BytesIO objects)

    Returns:
        Tuple of (page_count, confidence):
        - page_count: Estimated number of pages/images
        - confidence: 'high' (metadata), 'medium' (size heuristic), or 'low' (default fallback)

    Examples:
        >>> estimate_document_pages(Path("contract.pdf"))
        (23, 'high')  # Read from PDF metadata

        >>> estimate_document_pages(Path("scan.jpg"))
        (1, 'high')  # Single image

        >>> estimate_document_pages(io.BytesIO(pdf_bytes), "document.pdf")
        (15, 'medium')  # File size heuristic
    """
    try:
        # Determine file extension
        if hasattr(file_path, 'name'):
            # File-like object (BytesIO, UploadedFile)
            file_ext = Path(file_path.name).suffix.lower()
            file_size = _get_file_size(file_path)
        elif filename:
            # BytesIO with explicit filename
            file_ext = Path(filename).suffix.lower()
            file_size = _get_file_size(file_path)
        else:
            # Path object or string
            file_path = Path(file_path)
            file_ext = file_path.suffix.lower()
            file_size = file_path.stat().st_size

        # === PDF: Try pypdfium2 metadata read ===
        if file_ext == '.pdf':
            return _estimate_pdf_pages(file_path, file_size)

        # === DOCX/PPTX: Try python-docx/python-pptx ===
        elif file_ext in ['.docx', '.doc']:
            return _estimate_docx_pages(file_path, file_size)

        elif file_ext in ['.pptx', '.ppt']:
            return _estimate_pptx_slides(file_path, file_size)

        # === Images: Single page ===
        elif file_ext in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
            return (1, 'high')

        # === Text files: Estimate from size ===
        elif file_ext in ['.txt', '.md', '.html', '.xml']:
            # Rough estimate: 3000 chars per page
            pages = max(1, file_size // 3000)
            return (pages, 'medium')

        # === Fallback: Use default ===
        else:
            logger.warning(
                f"⚠️ Unknown file type '{file_ext}' - using default estimate "
                f"({DEFAULT_PAGE_ESTIMATE} pages)"
            )
            return (DEFAULT_PAGE_ESTIMATE, 'low')

    except Exception as e:
        logger.warning(f"⚠️ Page estimation failed: {e} - using default ({DEFAULT_PAGE_ESTIMATE})")
        return (DEFAULT_PAGE_ESTIMATE, 'low')


def _get_file_size(file_obj) -> int:
    """Get file size from file-like object or path"""
    if isinstance(file_obj, (Path, str)):
        return Path(file_obj).stat().st_size
    elif hasattr(file_obj, 'getbuffer'):
        # BytesIO
        return len(file_obj.getbuffer())
    elif hasattr(file_obj, 'seek') and hasattr(file_obj, 'tell'):
        # File-like object with seek/tell
        current_pos = file_obj.tell()
        file_obj.seek(0, 2)  # Seek to end
        size = file_obj.tell()
        file_obj.seek(current_pos)  # Restore position
        return size
    else:
        raise ValueError(f"Cannot determine file size for type: {type(file_obj)}")


def _estimate_pdf_pages(file_path, file_size: int) -> Tuple[int, str]:
    """
    Estimate PDF page count using pypdfium2 metadata.

    Falls back to size heuristic if library unavailable or read fails.
    """
    try:
        import pypdfium2 as pdfium

        # Handle both Path and BytesIO
        if isinstance(file_path, (Path, str)):
            pdf = pdfium.PdfDocument(str(file_path))
        elif hasattr(file_path, 'getbuffer'):
            # BytesIO - read bytes
            pdf_bytes = file_path.getvalue()
            pdf = pdfium.PdfDocument(pdf_bytes)
        elif hasattr(file_path, 'read'):
            # File-like object
            current_pos = file_path.tell()
            file_path.seek(0)
            pdf_bytes = file_path.read()
            file_path.seek(current_pos)
            pdf = pdfium.PdfDocument(pdf_bytes)
        else:
            raise ValueError(f"Unsupported file_path type for PDF: {type(file_path)}")

        page_count = len(pdf)
        pdf.close()

        logger.debug(f"✅ PDF page count from metadata: {page_count} pages")
        return (page_count, 'high')

    except ImportError:
        logger.debug("pypdfium2 not available - falling back to file size heuristic")
        return _estimate_from_file_size(file_size, SIZE_PER_PAGE_SCANNED_PDF)

    except Exception as e:
        logger.debug(f"PDF metadata read failed ({e}) - using file size heuristic")
        return _estimate_from_file_size(file_size, SIZE_PER_PAGE_SCANNED_PDF)


def _estimate_docx_pages(file_path, file_size: int) -> Tuple[int, str]:
    """
    Estimate DOCX page count (approximate - DOCX doesn't store page count directly).

    Falls back to size heuristic.
    """
    try:
        from docx import Document

        # python-docx only works with file paths or BytesIO
        if isinstance(file_path, (Path, str)):
            doc = Document(str(file_path))
        elif hasattr(file_path, 'read'):
            doc = Document(file_path)
        else:
            raise ValueError(f"Unsupported file_path type for DOCX: {type(file_path)}")

        # Rough estimate: ~3 paragraphs per page
        paragraph_count = len(doc.paragraphs)
        page_estimate = max(1, paragraph_count // 3)

        logger.debug(f"📄 DOCX page estimate from paragraph count: {page_estimate} pages")
        return (page_estimate, 'medium')

    except ImportError:
        logger.debug("python-docx not available - falling back to file size heuristic")
        return _estimate_from_file_size(file_size, SIZE_PER_PAGE_DOCX)

    except Exception as e:
        logger.debug(f"DOCX metadata read failed ({e}) - using file size heuristic")
        return _estimate_from_file_size(file_size, SIZE_PER_PAGE_DOCX)


def _estimate_pptx_slides(file_path, file_size: int) -> Tuple[int, str]:
    """
    Estimate PPTX slide count from metadata.

    Falls back to size heuristic.
    """
    try:
        from pptx import Presentation

        if isinstance(file_path, (Path, str)):
            prs = Presentation(str(file_path))
        elif hasattr(file_path, 'read'):
            prs = Presentation(file_path)
        else:
            raise ValueError(f"Unsupported file_path type for PPTX: {type(file_path)}")

        slide_count = len(prs.slides)

        logger.debug(f"📊 PPTX slide count from metadata: {slide_count} slides")
        return (slide_count, 'high')

    except ImportError:
        logger.debug("python-pptx not available - falling back to file size heuristic")
        return _estimate_from_file_size(file_size, SIZE_PER_PAGE_GENERIC)

    except Exception as e:
        logger.debug(f"PPTX metadata read failed ({e}) - using file size heuristic")
        return _estimate_from_file_size(file_size, SIZE_PER_PAGE_GENERIC)


def _estimate_from_file_size(file_size: int, bytes_per_page: int) -> Tuple[int, str]:
    """
    Fallback: Estimate page count from file size.

    Args:
        file_size: File size in bytes
        bytes_per_page: Average bytes per page for this file type

    Returns:
        (estimated_pages, 'medium')
    """
    pages = max(1, file_size // bytes_per_page)
    logger.debug(
        f"📏 File size estimate: {file_size:,} bytes ÷ {bytes_per_page:,} bytes/page "
        f"= {pages} pages"
    )
    return (pages, 'medium')


def get_confidence_message(confidence: str, page_count: int) -> str:
    """
    Get user-friendly confidence message for UI display.

    Args:
        confidence: 'high', 'medium', or 'low'
        page_count: Estimated page count

    Returns:
        Human-readable confidence message

    Examples:
        >>> get_confidence_message('high', 15)
        '~15 pages (from metadata)'

        >>> get_confidence_message('medium', 12)
        '~12 pages (estimated from file size ±30%)'

        >>> get_confidence_message('low', 15)
        '~15 pages (default estimate - actual may vary)'
    """
    if confidence == 'high':
        return f"~{page_count} pages (from metadata)"
    elif confidence == 'medium':
        return f"~{page_count} pages (estimated from file size ±30%)"
    else:  # low
        return f"~{page_count} pages (default estimate - actual may vary)"
